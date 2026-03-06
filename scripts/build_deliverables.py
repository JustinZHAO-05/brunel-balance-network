from __future__ import annotations

import json
import textwrap
from pathlib import Path
from typing import Any

import pandas as pd
from jinja2 import Template
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
ARTIFACTS = ROOT / "artifacts"
DELIVERABLES = ROOT / "deliverables"
ASSETS = DELIVERABLES / "assets"
SLIDE_ASSETS = ASSETS / "slides"

QUICK_DIR = ARTIFACTS / "reproductions" / "quick"
POSTER_DIR = ARTIFACTS / "reproductions" / "poster"
FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

PRIMARY = RGBColor(11, 27, 50)
ACCENT = RGBColor(31, 78, 121)
ORANGE = RGBColor(194, 65, 12)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(185, 28, 28)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(203, 213, 225)

PIL_PRIMARY = (11, 27, 50)
PIL_ACCENT = (31, 78, 121)
PIL_ORANGE = (194, 65, 12)
PIL_GREEN = (5, 150, 105)
PIL_RED = (185, 28, 28)
PIL_SLATE = (71, 85, 105)
PIL_LIGHT = (248, 250, 252)
PIL_WHITE = (255, 255, 255)
PIL_BORDER = (203, 213, 225)
PIL_PALE_BLUE = (239, 246, 255)
PIL_PALE_GREEN = (240, 253, 244)
PIL_PALE_ORANGE = (255, 247, 237)


def read_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def load_bundle(bundle_dir: Path) -> dict[str, Any]:
    panel_metrics = pd.read_csv(bundle_dir / "panel_metrics.csv")
    sweep_summary = pd.read_csv(bundle_dir / "sweep_summary.csv")
    acceptance = read_json(bundle_dir / "acceptance.json")
    return {
        "dir": bundle_dir,
        "panel_metrics": panel_metrics,
        "sweep_summary": sweep_summary,
        "acceptance": acceptance,
        "figure1": bundle_dir / "figure1_panels.png",
        "figure2": bundle_dir / "figure2_phase_diagram.png",
        "figure3": bundle_dir / "figure3_metrics.png",
    }


def format_float(value: Any, digits: int = 3) -> str:
    if pd.isna(value):
        return "NA"
    return f"{float(value):.{digits}f}"


def markdown_table(df: pd.DataFrame) -> str:
    headers = list(df.columns)
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]
    for _, row in df.iterrows():
        values = [str(row[col]).replace("|", "\\|") for col in headers]
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines)


def html_table(df: pd.DataFrame) -> str:
    return df.to_html(index=False, border=0, classes="data-table")


def build_summary_tables(quick: dict[str, Any], poster: dict[str, Any]) -> dict[str, pd.DataFrame]:
    quick_panels = quick["panel_metrics"].copy()
    poster_panels = poster["panel_metrics"].copy()
    for df in (quick_panels, poster_panels):
        for col in [
            "mean_firing_rate_hz",
            "median_cv_isi",
            "pairwise_correlation",
            "psd_peak_power",
        ]:
            df[col] = df[col].map(lambda value: format_float(value, 3))
    quick_table = quick_panels[["panel", "mean_firing_rate_hz", "median_cv_isi", "pairwise_correlation", "psd_peak_power"]]
    poster_table = poster_panels[["panel", "mean_firing_rate_hz", "median_cv_isi", "pairwise_correlation", "psd_peak_power"]]
    quick_display = quick_table.rename(
        columns={
            "mean_firing_rate_hz": "rate",
            "median_cv_isi": "CV",
            "pairwise_correlation": "corr",
            "psd_peak_power": "PSD",
        }
    )
    poster_display = poster_table.rename(
        columns={
            "mean_firing_rate_hz": "rate",
            "median_cv_isi": "CV",
            "pairwise_correlation": "corr",
            "psd_peak_power": "PSD",
        }
    )

    rule_rows = []
    for label, acceptance in [("quick", quick["acceptance"]), ("poster", poster["acceptance"])]:
        for rule_name, payload in acceptance["rules"].items():
            rule_rows.append(
                {
                    "preset": label,
                    "rule": rule_name,
                    "passed": "PASS" if payload["passed"] else "FAIL",
                    "detail": payload["detail"],
                }
            )
    acceptance_table = pd.DataFrame(rule_rows)
    acceptance_display = acceptance_table.rename(columns={"preset": "set", "passed": "status"})[["set", "rule", "status"]]
    acceptance_display["rule"] = acceptance_display["rule"].replace(
        {
            "async_corr_lt_sync": "async corr < sync",
            "async_cv_near_irregular": "async CV near irregular",
            "required_outputs_present": "outputs present",
            "sync_peak_power_gt_async": "sync peak > async",
        }
    )

    sweep = poster["sweep_summary"]
    heatmap_points = (
        sweep.sort_values("psd_peak_power", ascending=False)
        .head(6)[["g", "nu_ext_over_nu_thr", "mean_firing_rate_hz", "median_cv_isi", "psd_peak_power"]]
        .copy()
    )
    for col in ["g", "nu_ext_over_nu_thr", "mean_firing_rate_hz", "median_cv_isi", "psd_peak_power"]:
        heatmap_points[col] = heatmap_points[col].map(lambda value: format_float(value, 3))
    sweep_display = heatmap_points.head(4).rename(
        columns={
            "nu_ext_over_nu_thr": "nu",
            "mean_firing_rate_hz": "rate",
            "median_cv_isi": "CV",
            "psd_peak_power": "PSD",
        }
    )[["g", "nu", "rate", "PSD"]]

    return {
        "quick_panels": quick_table,
        "poster_panels": poster_table,
        "quick_panels_display": quick_display,
        "poster_panels_display": poster_display,
        "acceptance": acceptance_table,
        "acceptance_display": acceptance_display,
        "sweep_highlights": heatmap_points,
        "sweep_highlights_display": sweep_display,
    }


def write_report(quick: dict[str, Any], poster: dict[str, Any], tables: dict[str, pd.DataFrame]) -> None:
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    report_md = DELIVERABLES / "project_report.md"
    report_html = DELIVERABLES / "project_report.html"

    quick_acc = quick["acceptance"]
    poster_acc = poster["acceptance"]
    template = Template(
        textwrap.dedent(
            """
            # Brunel Balance Network 项目报告

            ## 1. 项目概述

            本项目实现了一个基于 Brian2 的 Brunel 平衡网络复现仓库，目标是用一个足够小、但足够真实的计算神经科学任务验证 `codex GPT-5.4 vibe coding` 的端到端交付能力。项目覆盖模型配置、仿真、分析、扫参、自动验收，以及最终的公开展示材料生成。

            ## 2. 研究问题

            本项目要回答的问题不是“能否精确重现论文中的每一个数值”，而是：

            > 在真实的 Brunel 平衡网络任务里，代理是否能端到端地完成可运行、可解释、可展示、可开源的研究型微项目？

            ## 3. 模型与实验设计

            - 模型：current-based LIF 兴奋/抑制网络
            - 框架：Brian2 2.7.1
            - 关键参数：抑制强度比 `g`、外部输入强度 `nu_ext_over_nu_thr`
            - 核心预设：A/B/C/D 四个 canonical panels
            - 扫参：在 `g x nu_ext_over_nu_thr` 网格上批量运行，构建小型相图

            ## 4. 复现工作流

            - `src.simulate`：单次仿真
            - `src.analyze`：单次结果分析与作图
            - `src.sweep`：批量扫参
            - `src.reproduce --preset quick|poster`：端到端复现

            ## 5. 主要图表

            ### Figure 1. Canonical panels

            ![Figure 1](../artifacts/reproductions/poster/figure1_panels.png)

            ### Figure 2. Poster-scale phase diagram

            ![Figure 2](../artifacts/reproductions/poster/figure2_phase_diagram.png)

            ### Figure 3. Metrics and acceptance

            ![Figure 3](../artifacts/reproductions/poster/figure3_metrics.png)

            ## 6. Quick 结果摘要

            {{ quick_table }}

            Quick 验收：`{{ quick_passed }}/{{ quick_total }}` 规则通过。

            ## 7. Poster 结果摘要

            {{ poster_table }}

            Poster 验收：`{{ poster_passed }}/{{ poster_total }}` 规则通过。

            ## 8. 验收规则逐项结果

            {{ acceptance_table }}

            ## 9. Poster Sweep 高亮点

            下表列出 `poster` 扫参中 `PSD peak power` 最高的 6 个点，用于快速识别更明显的同步态区域。

            {{ sweep_highlights }}

            ## 10. 结果解释

            - A/B 条件在 population-rate peak power 上显著高于 C/D，这与同步活动更强的预期一致。
            - C/D 条件在 pairwise correlation 的绝对值上整体低于 A/B，说明异步特征方向基本成立。
            - 当前项目的主要缺口是 `async_cv_near_irregular` 规则未通过。无论在 `quick` 还是 `poster` 预设下，C/D 的 `median CV(ISI)` 都低于目标阈值 `0.8`，说明网络缩放与参数保持之间仍有偏差。
            - 这类失败被保留在自动验收文件中，而不是被手工美化或隐藏，这一点本身也是 vibe coding 流程可信度的一部分。

            ## 11. 工程产出

            - 代码仓库：仿真、分析、扫参与复现 CLI
            - 数据产物：spike、population rate、metrics、phase diagram、acceptance JSON
            - 展示产物：项目报告、A1 poster、汇报 PPT
            - 过程记录：`docs/vibe_protocol.md` 与 `docs/vibe_log.md`

            ## 12. 局限与后续工作

            - 当前 `quick` 和 `poster` 都没有让 D 条件达到接近 1 的不规则放电水平
            - 网络缩放目前主要缩放 `N_E`，但对 Brunel 模型而言，如何保持连接数、输入强度和统计结构的一致性仍需更严谨处理
            - 下一步应优先做参数缩放校正，而不是继续堆更多图表

            ## 13. 复现命令

            ```bash
            python -m src.reproduce --preset quick
            python -m src.reproduce --preset poster
            python scripts/build_deliverables.py
            ```
            """
        ).strip()
    )

    markdown = template.render(
        quick_table=markdown_table(tables["quick_panels"]),
        poster_table=markdown_table(tables["poster_panels"]),
        acceptance_table=markdown_table(tables["acceptance"]),
        sweep_highlights=markdown_table(tables["sweep_highlights"]),
        quick_passed=quick_acc["passed_count"],
        quick_total=quick_acc["rule_count"],
        poster_passed=poster_acc["passed_count"],
        poster_total=poster_acc["rule_count"],
    )
    report_md.write_text(markdown + "\n", encoding="utf-8")

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <title>Brunel Balance Network 项目报告</title>
  <style>
    body {{ font-family: Arial, Helvetica, sans-serif; margin: 0; background: #f8fafc; color: #0f172a; }}
    main {{ max-width: 1080px; margin: 0 auto; padding: 48px 32px 80px; }}
    h1, h2, h3 {{ color: #0b1b32; }}
    p, li {{ line-height: 1.65; }}
    img {{ max-width: 100%; border: 1px solid #cbd5e1; border-radius: 12px; background: white; }}
    table {{ border-collapse: collapse; width: 100%; margin: 16px 0 32px; background: white; }}
    th, td {{ border: 1px solid #cbd5e1; padding: 10px 12px; text-align: left; font-size: 14px; }}
    th {{ background: #e2e8f0; }}
    code, pre {{ background: #e2e8f0; border-radius: 6px; }}
    pre {{ padding: 16px; overflow-x: auto; }}
    .hero {{ background: linear-gradient(135deg, #0b1b32, #1f4e79); color: white; padding: 28px 32px; border-radius: 18px; margin-bottom: 32px; }}
  </style>
</head>
<body>
  <main>
    <div class="hero">
      <h1>Brunel Balance Network 项目报告</h1>
      <p>Brian2-based computational neuroscience microproject built with codex GPT-5.4 vibe coding.</p>
    </div>
    <h2>1. 项目概述</h2>
    <p>本项目实现了一个基于 Brian2 的 Brunel 平衡网络复现仓库，目标是用一个足够小、但足够真实的计算神经科学任务验证 <code>codex GPT-5.4 vibe coding</code> 的端到端交付能力。项目覆盖模型配置、仿真、分析、扫参、自动验收，以及最终的公开展示材料生成。</p>
    <h2>2. 研究问题</h2>
    <p>在真实的 Brunel 平衡网络任务里，代理是否能端到端地完成可运行、可解释、可展示、可开源的研究型微项目？</p>
    <h2>3. 模型与实验设计</h2>
    <ul>
      <li>模型：current-based LIF 兴奋/抑制网络</li>
      <li>框架：Brian2 2.7.1</li>
      <li>关键参数：<code>g</code>、<code>nu_ext / nu_thr</code></li>
      <li>核心预设：A/B/C/D 四个 canonical panels</li>
      <li>扫参：在 <code>g x nu_ext / nu_thr</code> 网格上批量运行</li>
    </ul>
    <h2>4. 主要图表</h2>
    <h3>Figure 1. Canonical panels</h3>
    <img src="../artifacts/reproductions/poster/figure1_panels.png" alt="Figure 1" />
    <h3>Figure 2. Poster-scale phase diagram</h3>
    <img src="../artifacts/reproductions/poster/figure2_phase_diagram.png" alt="Figure 2" />
    <h3>Figure 3. Metrics and acceptance</h3>
    <img src="../artifacts/reproductions/poster/figure3_metrics.png" alt="Figure 3" />
    <h2>5. Quick 结果摘要</h2>
    {html_table(tables["quick_panels"])}
    <p>Quick 验收：<strong>{quick_acc["passed_count"]}/{quick_acc["rule_count"]}</strong> 规则通过。</p>
    <h2>6. Poster 结果摘要</h2>
    {html_table(tables["poster_panels"])}
    <p>Poster 验收：<strong>{poster_acc["passed_count"]}/{poster_acc["rule_count"]}</strong> 规则通过。</p>
    <h2>7. 验收规则逐项结果</h2>
    {html_table(tables["acceptance"])}
    <h2>8. Poster Sweep 高亮点</h2>
    {html_table(tables["sweep_highlights"])}
    <h2>9. 结果解释</h2>
    <ul>
      <li>A/B 条件在 population-rate peak power 上显著高于 C/D，与同步活动更强的预期一致。</li>
      <li>C/D 条件在 pairwise correlation 的绝对值上整体低于 A/B，异步特征方向基本成立。</li>
      <li>当前主要缺口是 <code>async_cv_near_irregular</code> 规则未通过，说明 C/D 的不规则性仍不够强。</li>
      <li>这一失败被保留在自动验收中，而不是被手工美化或隐藏。</li>
    </ul>
    <h2>10. 复现命令</h2>
    <pre><code>python -m src.reproduce --preset quick
python -m src.reproduce --preset poster
python scripts/build_deliverables.py</code></pre>
  </main>
</body>
</html>"""
    report_html.write_text(html, encoding="utf-8")


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, size=size)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split(" ")
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = word
    if current:
        lines.append(current)
    if not lines:
        lines = [text]
    return lines


def draw_multiline(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    max_width: int,
    fill: tuple[int, int, int],
    line_gap: int = 8,
    bullet: bool = False,
) -> int:
    x, y = xy
    lines = wrap_text(draw, text, font, max_width - (26 if bullet else 0))
    bullet_text = "• " if bullet else ""
    for idx, line in enumerate(lines):
        prefix = bullet_text if idx == 0 else "  " if bullet else ""
        draw.text((x, y), prefix + line, font=font, fill=fill)
        bbox = draw.textbbox((x, y), prefix + line, font=font)
        y += (bbox[3] - bbox[1]) + line_gap
    return y


def draw_card_pil(
    image: Image.Image,
    box: tuple[int, int, int, int],
    title: str,
    body_lines: list[str],
    fill: tuple[int, int, int] = PIL_WHITE,
) -> None:
    draw = ImageDraw.Draw(image)
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=28, fill=fill, outline=PIL_BORDER, width=3)
    draw.text((x1 + 28, y1 + 22), title, font=pil_font(34, bold=True), fill=PIL_PRIMARY)
    y = y1 + 82
    for line in body_lines:
        y = draw_multiline(draw, (x1 + 28, y), line, pil_font(23), x2 - x1 - 56, PIL_SLATE, line_gap=8, bullet=True)
        y += 6


def paste_contained(base: Image.Image, source_path: Path, box: tuple[int, int, int, int], padding: int = 0, bg=PIL_WHITE) -> None:
    x1, y1, x2, y2 = box
    target_w = max(1, x2 - x1 - 2 * padding)
    target_h = max(1, y2 - y1 - 2 * padding)
    canvas = Image.new("RGB", (x2 - x1, y2 - y1), bg)
    source = Image.open(source_path).convert("RGB")
    fitted = ImageOps.contain(source, (target_w, target_h))
    paste_x = (canvas.width - fitted.width) // 2
    paste_y = (canvas.height - fitted.height) // 2
    canvas.paste(fitted, (paste_x, paste_y))
    base.paste(canvas, (x1, y1))


def paste_contained_crop(
    base: Image.Image,
    source_path: Path,
    box: tuple[int, int, int, int],
    crop: tuple[float, float, float, float],
    padding: int = 0,
    bg=PIL_WHITE,
) -> None:
    source = Image.open(source_path).convert("RGB")
    width, height = source.size
    left = int(width * crop[0])
    top = int(height * crop[1])
    right = int(width * crop[2])
    bottom = int(height * crop[3])
    cropped = source.crop((left, top, right, bottom))
    temp = ASSETS / "_temp_crop.png"
    cropped.save(temp)
    paste_contained(base, temp, box, padding=padding, bg=bg)
    temp.unlink(missing_ok=True)


def draw_dataframe(
    image: Image.Image,
    df: pd.DataFrame,
    box: tuple[int, int, int, int],
    title: str,
    header_fill: tuple[int, int, int] = PIL_PALE_BLUE,
    body_fill: tuple[int, int, int] = PIL_WHITE,
) -> None:
    draw = ImageDraw.Draw(image)
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=24, fill=body_fill, outline=PIL_BORDER, width=3)
    draw.text((x1 + 24, y1 + 18), title, font=pil_font(28, bold=True), fill=PIL_PRIMARY)
    table_top = y1 + 70
    cols = len(df.columns)
    rows = len(df) + 1
    col_w = (x2 - x1 - 40) / cols
    row_h = (y2 - table_top - 20) / rows
    for col_idx, column in enumerate(df.columns):
        cell = (
            x1 + 20 + int(col_idx * col_w),
            table_top,
            x1 + 20 + int((col_idx + 1) * col_w),
            int(table_top + row_h),
        )
        draw.rectangle(cell, fill=header_fill, outline=PIL_BORDER, width=2)
        draw.text((cell[0] + 10, cell[1] + 8), str(column), font=pil_font(17, bold=True), fill=PIL_PRIMARY)
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, column in enumerate(df.columns):
            cell = (
                x1 + 20 + int(col_idx * col_w),
                int(table_top + row_idx * row_h),
                x1 + 20 + int((col_idx + 1) * col_w),
                int(table_top + (row_idx + 1) * row_h),
            )
            draw.rectangle(cell, fill=PIL_WHITE, outline=PIL_BORDER, width=1)
            draw.text((cell[0] + 10, cell[1] + 8), str(row[column]), font=pil_font(16), fill=PIL_SLATE)


def create_poster_png(quick: dict[str, Any], poster: dict[str, Any], tables: dict[str, pd.DataFrame]) -> Path:
    width, height = 4967, 3508
    image = Image.new("RGB", (width, height), PIL_LIGHT)
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, width, 260), radius=0, fill=PIL_PRIMARY)
    draw.text((100, 70), "Brunel Balance Network", font=pil_font(74, bold=True), fill=PIL_WHITE)
    draw.text((100, 155), "Brian2-based computational neuroscience microproject built with codex GPT-5.4 vibe coding", font=pil_font(28), fill=(226, 232, 240))
    draw.rounded_rectangle((4000, 56, 4820, 190), radius=38, fill=PIL_ORANGE)
    draw.text((4060, 95), "A1 Poster", font=pil_font(34, bold=True), fill=PIL_WHITE)

    draw_card_pil(
        image,
        (80, 320, 960, 1060),
        "Project question",
        [
            "Can an agent complete a real but bounded computational neuroscience workflow end-to-end?",
            "The task includes simulation, analysis, sweep experiments, acceptance checks, and open-source packaging.",
        ],
    )
    draw_card_pil(
        image,
        (80, 1100, 960, 1820),
        "Model and setup",
        [
            "Current-based LIF E/I network in Brian2 2.7.1",
            "Canonical A/B/C/D panels adapted from the Brian2 Brunel_2000 example",
            "Key axes: inhibitory strength ratio g and external drive nu_ext / nu_thr",
        ],
        fill=PIL_PALE_BLUE,
    )
    draw_card_pil(
        image,
        (80, 1860, 960, 2760),
        "Acceptance snapshot",
        [
            f"Quick preset: {quick['acceptance']['passed_count']}/{quick['acceptance']['rule_count']} rules passed",
            f"Poster preset: {poster['acceptance']['passed_count']}/{poster['acceptance']['rule_count']} rules passed",
            "Stable wins: sync peak power, output completeness, lower async correlation",
            "Open issue: async CV(ISI) remains below the target irregularity threshold",
        ],
        fill=PIL_PALE_ORANGE,
    )
    draw_card_pil(
        image,
        (80, 2800, 960, 3410),
        "Repository",
        [
            "github.com/JustinZHAO-05/brunel-balance-network",
            "Deliverables in this repo: report, poster, presentation PPT",
            "The acceptance failure is kept visible instead of being hand-corrected.",
        ],
        fill=PIL_PALE_GREEN,
    )

    paste_contained(image, poster["figure1"], (1020, 320, 2885, 2000), padding=18)
    draw.text((1040, 2015), "Figure 1. Poster-scale canonical panels", font=pil_font(28, bold=True), fill=PIL_PRIMARY)
    paste_contained_crop(image, poster["figure3"], (1020, 2080, 2885, 3410), crop=(0.0, 0.0, 1.0, 0.86), padding=18)
    draw.text((1040, 3430 - 24), "Figure 3. Metrics and automated acceptance", font=pil_font(28, bold=True), fill=PIL_PRIMARY)

    paste_contained(image, poster["figure2"], (2945, 320, 4885, 1680), padding=18)
    draw.text((2970, 1696), "Figure 2. Poster sweep phase diagram", font=pil_font(28, bold=True), fill=PIL_PRIMARY)
    draw_dataframe(image, tables["poster_panels_display"], (2945, 1760, 4885, 2440), "Poster panel metrics")
    draw_dataframe(image, tables["sweep_highlights_display"], (2945, 2490, 4885, 3410), "Top PSD peak power points")

    out_path = ASSETS / "poster_a1.png"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(out_path, quality=95)
    return out_path


def create_slide_pngs(quick: dict[str, Any], poster: dict[str, Any], tables: dict[str, pd.DataFrame]) -> list[Path]:
    SLIDE_ASSETS.mkdir(parents=True, exist_ok=True)

    def new_slide() -> Image.Image:
        return Image.new("RGB", (1600, 900), PIL_LIGHT)

    slide_paths: list[Path] = []

    # Slide 1
    img = new_slide()
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((0, 0, 1600, 900), radius=0, fill=PIL_LIGHT)
    draw.rounded_rectangle((0, 0, 1600, 900), radius=0, fill=PIL_LIGHT)
    draw.rounded_rectangle((40, 40, 660, 860), radius=36, fill=PIL_WHITE, outline=PIL_BORDER, width=3)
    draw.text((90, 120), "Brunel Balance Network", font=pil_font(46, bold=True), fill=PIL_PRIMARY)
    y = 220
    for line in [
        "Brian2-based computational neuroscience microproject",
        "Built with codex GPT-5.4 vibe coding",
        "Goal: test whether an agent can deliver a research-style repo end-to-end",
    ]:
        y = draw_multiline(draw, (95, y), line, pil_font(24), 500, PIL_SLATE)
        y += 12
    draw.rounded_rectangle((90, 520, 420, 600), radius=28, fill=PIL_ORANGE)
    draw.text((125, 545), "7-slide briefing", font=pil_font(24, bold=True), fill=PIL_WHITE)
    paste_contained(img, poster["figure1"], (720, 70, 1540, 830), padding=12)
    p = SLIDE_ASSETS / "slide_01_title.png"
    img.save(p); slide_paths.append(p)

    # Slide 2
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.text((70, 60), "Problem and scope", font=pil_font(40, bold=True), fill=PIL_PRIMARY)
    draw.text((70, 115), "What this project tries to validate", font=pil_font(20), fill=PIL_SLATE)
    draw_card_pil(img, (60, 180, 500, 760), "Question", [
        "Can an agent complete a small but real computational neuroscience project from code to public presentation?"
    ])
    draw_card_pil(img, (540, 180, 980, 760), "Scope", [
        "Model implementation",
        "Configuration presets and sweeps",
        "Analysis metrics and plots",
        "Acceptance rules and repo packaging",
    ], fill=PIL_PALE_BLUE)
    draw_card_pil(img, (1020, 180, 1540, 760), "Guardrail", [
        "Failures remain visible. The repo does not pretend every rule passed."
    ], fill=PIL_PALE_ORANGE)
    p = SLIDE_ASSETS / "slide_02_scope.png"
    img.save(p); slide_paths.append(p)

    # Slide 3
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.text((70, 60), "Workflow and CLI", font=pil_font(40, bold=True), fill=PIL_PRIMARY)
    draw.text((70, 115), "Four layers, each responsible for one part of the pipeline", font=pil_font(20), fill=PIL_SLATE)
    cli_df = pd.DataFrame(
        [
            ["simulate", "single run", "spikes.npz, population_rate.csv"],
            ["analyze", "single-run metrics", "metrics.json, raster.png, summary.png"],
            ["sweep", "parameter grid", "summary.csv, phase_diagram.png"],
            ["reproduce", "end-to-end workflow", "Figure 1/2/3, acceptance.json"],
        ],
        columns=["CLI", "Role", "Outputs"],
    )
    draw_dataframe(img, cli_df, (70, 180, 1530, 530), "CLI map")
    draw_card_pil(img, (70, 570, 700, 820), "Recommended run order", [
        "1. simulate single panel",
        "2. analyze that run",
        "3. reproduce --preset quick",
        "4. reproduce --preset poster",
    ], fill=PIL_PALE_GREEN)
    draw_card_pil(img, (760, 570, 1530, 820), "Reasoning", [
        "The split keeps simulation, analysis, and presentation artifacts independently reproducible."
    ])
    p = SLIDE_ASSETS / "slide_03_workflow.png"
    img.save(p); slide_paths.append(p)

    # Slide 4
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.text((70, 60), "Canonical panels", font=pil_font(40, bold=True), fill=PIL_PRIMARY)
    draw.text((70, 115), "A/B/C/D poster preset outputs", font=pil_font(20), fill=PIL_SLATE)
    paste_contained(img, poster["figure1"], (60, 160, 1160, 830), padding=12)
    draw_dataframe(img, tables["poster_panels_display"], (1190, 170, 1540, 540), "Poster panel metrics")
    draw_card_pil(img, (1190, 570, 1540, 830), "Readout", [
        "A/B show stronger rhythmic population activity than C/D.",
        "D remains low-rate and weakly irregular under current scaling.",
    ], fill=PIL_PALE_ORANGE)
    p = SLIDE_ASSETS / "slide_04_panels.png"
    img.save(p); slide_paths.append(p)

    # Slide 5
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.text((70, 60), "Parameter sweep", font=pil_font(40, bold=True), fill=PIL_PRIMARY)
    draw.text((70, 115), "Poster preset phase diagram and top synchronous points", font=pil_font(20), fill=PIL_SLATE)
    paste_contained(img, poster["figure2"], (60, 160, 980, 830), padding=12)
    draw_dataframe(img, tables["sweep_highlights_display"], (1020, 170, 1540, 600), "Highest PSD peak power")
    draw_card_pil(img, (1020, 630, 1540, 830), "Interpretation", [
        "High external drive and moderate-to-high inhibition produce the strongest oscillatory signatures in this sweep."
    ], fill=PIL_PALE_BLUE)
    p = SLIDE_ASSETS / "slide_05_sweep.png"
    img.save(p); slide_paths.append(p)

    # Slide 6
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.text((70, 60), "Acceptance and limitations", font=pil_font(40, bold=True), fill=PIL_PRIMARY)
    draw.text((70, 115), "The repo keeps the incomplete part visible", font=pil_font(20), fill=PIL_SLATE)
    paste_contained_crop(img, poster["figure3"], (60, 160, 900, 830), crop=(0.0, 0.0, 1.0, 0.86), padding=12)
    draw_dataframe(img, tables["acceptance_display"], (930, 170, 1540, 530), "Rule-by-rule status")
    draw_card_pil(img, (930, 600, 1540, 830), "Main limitation", [
        "Both quick and poster presets fail the async CV(ISI) threshold. The next technical step is scaling correction, not more window dressing."
    ], fill=PIL_PALE_ORANGE)
    p = SLIDE_ASSETS / "slide_06_acceptance.png"
    img.save(p); slide_paths.append(p)

    # Slide 7
    img = new_slide(); draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((0, 0, 1600, 900), radius=0, fill=PIL_PRIMARY)
    draw.text((70, 80), "Takeaways", font=pil_font(44, bold=True), fill=PIL_WHITE)
    draw.text((70, 145), "What this microproject demonstrates", font=pil_font(22), fill=(226, 232, 240))
    draw_card_pil(img, (70, 220, 520, 720), "Delivered", [
        "Complete repo",
        "Simulation and analysis pipeline",
        "Sweep and acceptance system",
        "Report, A1 poster, and PPT",
    ], fill=PIL_WHITE)
    draw_card_pil(img, (570, 220, 1030, 720), "Still weak", [
        "Irregular firing target in C/D",
        "Scaling consistency under reduced N_E",
        "Need stronger parameter-preserving reductions",
    ], fill=PIL_PALE_ORANGE)
    draw_card_pil(img, (1080, 220, 1530, 720), "Repo", [
        "github.com/JustinZHAO-05/brunel-balance-network",
        "Built with codex GPT-5.4 vibe coding",
        "Artifacts generated directly from simulation outputs",
    ], fill=PIL_PALE_BLUE)
    p = SLIDE_ASSETS / "slide_07_takeaways.png"
    img.save(p); slide_paths.append(p)

    return slide_paths


def set_textbox_text(shape, text: str, font_size: int, color: RGBColor = PRIMARY, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"


def add_title(slide, left, top, width, height, text, size=24):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_textbox_text(box, text, size, color=PRIMARY, bold=True)
    return box


def add_body(slide, left, top, width, height, lines, size=14, color=SLATE):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
    return box


def add_card(slide, left, top, width, height, title, body_lines, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    add_title(slide, left + Cm(0.35), top + Cm(0.2), width - Cm(0.7), Cm(0.9), title, size=18)
    add_body(slide, left + Cm(0.35), top + Cm(1.0), width - Cm(0.7), height - Cm(1.2), body_lines, size=12)
    return shape


def add_table(slide, left, top, width, height, df: pd.DataFrame, font_size: int = 11) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.first_row = True
    for col_idx, column in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.name = "Aptos"
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, column in enumerate(df.columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(row[column])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = "Aptos"


def build_poster_pptx(quick: dict[str, Any], poster: dict[str, Any], tables: dict[str, pd.DataFrame]) -> None:
    poster_png = create_poster_png(quick, poster, tables)
    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(23.39)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(poster_png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    prs.save(DELIVERABLES / "brunel_balance_network_poster_a1.pptx")


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    add_title(slide, Cm(0.9), Cm(0.6), Cm(22), Cm(1.1), title, size=24)
    if subtitle:
        add_body(slide, Cm(0.9), Cm(1.6), Cm(22), Cm(0.9), [subtitle], size=12)


def build_presentation_pptx(quick: dict[str, Any], poster: dict[str, Any], tables: dict[str, pd.DataFrame]) -> None:
    slide_pngs = create_slide_pngs(quick, poster, tables)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_path in slide_pngs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(slide_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(DELIVERABLES / "brunel_balance_network_presentation.pptx")


def write_deliverables_readme() -> None:
    text = """# Deliverables

- `project_report.md`: 项目长报告，适合仓库内阅读和继续修改
- `project_report.html`: 带样式的可直接打开报告
- `brunel_balance_network_poster_a1.pptx`: A1 尺寸单页 poster
- `brunel_balance_network_presentation.pptx`: 汇报 PPT
"""
    (DELIVERABLES / "README.md").write_text(text, encoding="utf-8")


def main() -> None:
    quick = load_bundle(QUICK_DIR)
    poster = load_bundle(POSTER_DIR)
    tables = build_summary_tables(quick, poster)

    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)

    write_report(quick, poster, tables)
    build_poster_pptx(quick, poster, tables)
    build_presentation_pptx(quick, poster, tables)
    write_deliverables_readme()


if __name__ == "__main__":
    main()
